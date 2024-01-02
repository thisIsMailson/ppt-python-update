from pptx import Presentation
# import win32com.client
from pptx.util import Pt
import requests
from bs4 import BeautifulSoup
import base64
from io import BytesIO
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
def make_request(url, headers):
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.json()['data']['attributes']
    else:
        raise Exception(f"Request failed with status code {response.status_code}")
    
headers = {
    'X-Api-Version': '3',
    'Authorization': API_TOKEN,
    'Accept': 'application/json'
}

def getImages(data):
    urls = []
    for item in data:
        if item['attributes'] and item['attributes'].get('image'):
            image_data = item['attributes'].get('image')
            upload_id = image_data['upload_id']
            try:
                response = requests.get(f'https://site-api.datocms.com/uploads/{upload_id}', headers=headers)
                response.raise_for_status()
                data = response.json()['data']['attributes']
                urls.append((data['url'], data['copyright']))
            except requests.exceptions.RequestException as e:
                raise Exception(f"Request failed with error: {e}")
        else:
            # Handle the case where 'image' key is not present in the item's attributes
            pass
    return urls

def image_to_base64(urls):
    images = []

    for url in urls:
        response = requests.get(url[0])
        image_content = response.content
        base64_image = base64.b64encode(image_content).decode('utf-8')

        image_content = BytesIO(response.content)
        images.append((image_content, url[1]))
    return images

def get_main_advantage(id):
    url = f'https://site-api.datocms.com/items/{id}?nested=true&version=published'
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.json()['data']['attributes']['name']

def get_advantages(data):
    results = []
    for item in data:
        item_id = item['id']
        url = f'https://site-api.datocms.com/items/{item_id}?nested=true&version=published'
        response = requests.get(url, headers=headers)
        if response.status_code == 200:
            advantage_id = response.json()['data']['attributes']['advantage']
            advantage = response.json()['data']['attributes']['value']
            main_advantage = get_main_advantage(advantage_id)
            results.append((main_advantage, advantage))
        else:
            results.append((None, None))
    return results

def move_content(elements_to_move, shapes, normalizer=1):
    gap_reduction = 0.4 * abs(normalizer-5)
    for element in elements_to_move:
        shapes[element].top -= Inches(gap_reduction)

def remove_cells(cells_to_remove, shapes):
    for index in cells_to_remove:
        cell = shapes[index]
        shapes._spTree.remove(cell._element)

def push_elements_down(elements, shapes, normalizer=1):
    for element in elements:
        shapes[element].top += Inches(0.12)
def generate_pptx():
    prod_id = 120784088

    url = f'https://site-api.datocms.com/items/{prod_id}?nested=true&version=published'
    try:
        product = make_request(url, headers)
       
        platform = product['platform']
        seg_url = f'https://site-api.datocms.com/items/{platform}?nested=true&version=published'
        segment = make_request(seg_url, headers)['name']
        cat = product['category']
        cat_url = f'https://site-api.datocms.com/items/{cat}?nested=true&version=published'
        category = make_request(cat_url, headers)['name']
        cont = product['contact_info']
        contact_url = f'https://site-api.datocms.com/items/{cont}?nested=true&version=published'
        contact = make_request(contact_url, headers)
        comp = product['company']
        company_url =  f'https://site-api.datocms.com/items/{comp}?nested=true&version=published'
        company = make_request(company_url, headers)
        div = product['division']
        division_url = f'https://site-api.datocms.com/items/{div}?nested=true&version=published'
        division = make_request(division_url, headers)
        adv_url = f'https://site-api.datocms.com/items/56982708?nested=true&version=published'
        adv = make_request(adv_url, headers)
        prod_trl = product['trl_level']
        trl_url = f'https://site-api.datocms.com/items/{prod_trl}?nested=true&version=published'
        tr_level = make_request(trl_url, headers)
        
        get_advantages(product['main_advantages'])
        
        advantage_tuple = get_advantages(product['main_advantages'])
        ###***********  insert production code for bucket here ***********###

        # bucket_name = 'powepoint_templates'
        # file_name = 'Motherson-Fact sheet_'+str(len(advantage_tuple))+'_rows.pptx'

        # storage_client = storage.Client()
        # bucket = storage_client.bucket(bucket_name)
        # blob = bucket.blob(file_name)


        # file_contents = blob.download_as_bytes()

        # factsheet_ppt = Presentation(io.BytesIO(file_contents))

        ####***********  local test code ***********###
        factsheet_ppt = Presentation('./templates/Motherson-Fact sheet_5_rows.pptx')
  

        slide = factsheet_ppt.slides[0]
        shapes = slide.shapes
        # for i in shapes:
        #     shapes[i].text = f'{i}'
        print(len(advantage_tuple))
        slide_shapes = [
            ('title', shapes[4]), #0
            ('segment', shapes[33]), #1
            ('category', shapes[32]), #2
            ('description', shapes[5]), #3
            ('key_facts', shapes[6]), #4
            ('applications_compliancy', shapes[7]), #5
            ('intellectual_property', shapes[8]), #6
            ('first_image', shapes[1]), #7
            ('second_image', shapes[2]), #8    
            ('third_image', shapes[3]), #9
            ('img_source_1', shapes[24]), #10
            ('img_source_2', shapes[26]), #11
            ('img_source_3', shapes[25]), #12
            ('contact', shapes[21]), #13
            ('division', shapes[22]), #14
            ('email', shapes[23]), #15
            ('table', shapes[30]), #16
            ('tr_bar_1', shapes[35]), #17
            ('tr_4', shapes[50]), #18 // this one is the text above the bar
            ('tr_bar_2', shapes[48]), #19
            ('tr_bar_3', shapes[49]), #20
            ('tr_bar_4', shapes[50]), #21
            ('intelectual_title', shapes[29]), #22
            ('key_facts_title', shapes[23]), #23
            ('applications_compliancy_title', shapes[24]), #24
            ('intellectual_property_title', shapes[25]), #25
        ]
        
        # for i in range(1, len(shapes)):
        num_characters = product['applications_compliancy'].count('\n') * 45
        num_paragraphs = product['applications_compliancy'].count('\n') + 1 # it starts at 0
        # Constants from the next calculations
        a = 0.002
        b = 0.04

        # Calculate the top shift based on the total number of characters using the formula
        top_shift = a * num_characters + b
        if top_shift < 0.1:
            top_shift = 0.0028 * len(product['applications_compliancy'])
        if num_paragraphs > 1:
            slide_shapes[6][1].top += Inches(top_shift)
            slide_shapes[22][1].top += Inches(top_shift)

        # Change the title
        title = slide_shapes[0][1].text_frame.paragraphs[0]
        t_run = title.runs[0]
        t_run.level = 0
        t_run.font.color.rgb = RGBColor(255, 0, 0) 
        title.text = BeautifulSoup(product['name'], "html.parser").get_text()

        # Change the segment
        sl_segment = slide_shapes[1][1].text_frame.paragraphs[0]
        run = sl_segment.runs[0]
        run.font.size = Pt(8)
        run.level = 0
        run.text = 'Segment: '+BeautifulSoup(segment, "html.parser").get_text()

        # Change the category
        sl_category = slide_shapes[2][1].text_frame.paragraphs[0]
        run = sl_category.runs[0]
        run.level = 0
        run.text = 'Category: '+ BeautifulSoup(category, "html.parser").get_text()
        run.font.size = Pt(8)

        # Change the description

        shapes_to_move_by_description = [6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16, 17, 18, 19, 20, 27, 28, 29, 30, 51, 52, 53, 54, 55, 56, 57, 58]
        if (len(product['description']) > 100):
            push_elements_down(shapes_to_move_by_description, shapes, 1)


        description = slide_shapes[3][1].text_frame.paragraphs[0]
        run = description.runs[0]
        run.font.size = Pt(8)
        run.level = 0
        run.text = BeautifulSoup(product['description'], "html.parser").get_text()

        # Change the key facts
        text_frame = slide_shapes[4][1].text_frame  # Text frame for 'key_facts'

        paragraphs = product['key_facts'].split('\n')

        # Handle the first paragraph separately
        if len(paragraphs) > 0:
            first_paragraph = text_frame.paragraphs[0]
            first_paragraph.clear()
            first_paragraph.text = paragraphs[0]
            first_paragraph.font.size = Pt(7)
            first_paragraph.alignment = PP_ALIGN.LEFT

            # Manually add bullet format for the first paragraph
            for run in first_paragraph.runs:
                run.text = '• ' + run.text

        # Add bullet points for the remaining paragraphs starting from the second paragraph
        for paragraph_text in paragraphs[1:]:
            p = text_frame.add_paragraph()
            p.text = '• ' + paragraph_text
            p.font.size = Pt(7)
            p.alignment = PP_ALIGN.LEFT

        # Change the key applicatiin and compliancy
        sl_application = slide_shapes[5][1].text_frame.paragraphs[0]
        run = sl_application.runs[0]
        run.font.size = Pt(7)
        run.level = 0
        run.text = BeautifulSoup(product['applications_compliancy'], "html.parser").get_text()

        # Change the intelectual property
        sl_intelect =slide_shapes[6][1].text_frame.paragraphs[0]
        run = sl_intelect.runs[0]
        run.font.size = Pt(7)
        run.level = 0
        run.text = product['intellectual_property']

        # Change the contact
        sl_contact = slide_shapes[13][1].text_frame.paragraphs[0]
        run = sl_contact.runs[0]
        run.font.size = Pt(6.3)
        run.level = 0
        c_text = 'Contact: ' + contact['name'] +'\n' +'Function: '+ contact['job_title']
        run.text =  BeautifulSoup(c_text, "html.parser").get_text()

        # Change the division
        sl_division = slide_shapes[14][1].text_frame.paragraphs[0]
        run = sl_division.runs[0]
        run.font.size = Pt(6.3)
        run.level = 0
        d_text = 'Division: ' + division['name'] +'\n' +'Company: '+ company['name']
        run.text =  BeautifulSoup(d_text, "html.parser").get_text()

        # Change the Email
        date_str = str(product['updated_at'])
        dt = datetime.fromisoformat(date_str)

        formatted_date = dt.strftime("%d-%m-%Y")

        sl_email = slide_shapes[15][1].text_frame.paragraphs[0]
        run = sl_email.runs[0]
        run.font.size = Pt(6.3)
        run.level = 0
        d_text = ''+contact['email'] +'\n' +'Last modified: '+ formatted_date
        run.text =  BeautifulSoup(d_text, "html.parser").get_text()

        # Change the table
        # deal with the cells
        cells_to_remove = {
            1: [11, 11, 11, 11, 49, 47, 48, 47, 49, 47, 48, 47], # cells are repeated since the index shifts after removing a cell
            2: [13, 13, 49, 51, 49, 50, 50, 50],
            3: [51, 51, 54, 54, 53],
            4: [55, 55]
        }
        elements_to_move = {
            1: [6, 7, 8, 27, 28, 29],
            2: [6, 7, 8, 27, 28, 29],
            3: [6, 7, 8, 27, 28, 29],
            4: [6, 7, 8, 27, 28, 29]
        }
        advantage_length = len(advantage_tuple)
        if advantage_length != 5:
            move_content(elements_to_move[advantage_length], shapes, advantage_length)

        if advantage_length in cells_to_remove:
            remove_cells(cells_to_remove[advantage_length], shapes)
        elif advantage_length == 5:
            pass

        # the following represents the cell index in each row that needs to be populated
        # so for example, if the advantage_tuple has 2 elements, we know that the two cells' indexes are 9 and 10 and so on and so forth
        # the cells index should be stored in the slide_shapes... TODO
        table_rows = {
            1: [9, 10],
            2: [9, 10, 11, 12],
            3: [9, 10, 11, 12, 13, 14],
            4: [9, 10, 11, 12, 13, 14, 51, 52],
            5: [9, 10, 11, 12, 13, 14, 51, 52, 55, 56] # the this last row we account for all the previous cells
        }

        # for i in range(1, len(shapes)):
        #     shapes[i].text = f'shape {i}'

        cells_to_populate = table_rows.get(len(advantage_tuple), [])
        for i, cell in enumerate(cells_to_populate):
            cell_index = i // 2  # Integer division to get the corresponding tuple index
            curr_cell = shapes[cell].text_frame.paragraphs[0]
            run = curr_cell.runs[0]
            run.font.size = Pt(6.3)
            run.level = 0
            if cell_index < len(advantage_tuple):
                if i % 2 == 0:  # Using i to determine if it's the first or second element of the tuple(cell wise)
                    run.text = BeautifulSoup(advantage_tuple[cell_index][0], "html.parser").get_text() # First element of the tuple for column 1
                else:
                    run.text = BeautifulSoup(advantage_tuple[cell_index][1], "html.parser").get_text()  # Second element of the tuple for column 2

        images = image_to_base64(getImages(product['attachments']))
        for i, (img_content, _) in enumerate(images):
            if i + 7 >= len(slide_shapes):
                break

            auto_shape = slide_shapes[i + 7][1]
            if not auto_shape.has_text_frame:
                continue

            left = auto_shape.left
            top = auto_shape.top
            width = auto_shape.width
            height = auto_shape.height
            
            slide.shapes._spTree.remove(auto_shape._element)

            new_picture = slide.shapes.add_picture(img_content, left, top, width, height)

            slide.shapes._spTree.remove(new_picture._element)
            slide.shapes._spTree.insert(2, new_picture._element)

        print( len(images) )

        # Handle the image sources
        if len(images) > 0:
            ref_element = slide_shapes[10][1]
            paragraph = ref_element.text_frame.paragraphs[0]
            run = paragraph.runs[0]
            run.text = images[0][1] or 'Source: Motherson'
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.size = Pt(8)
        else:
            shapes._spTree.remove(slide_shapes[10][1]._element)

        if len(images) > 1:
            ref_element = slide_shapes[11][1]
            paragraph = ref_element.text_frame.paragraphs[0]
            run = paragraph.runs[0]
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.size = Pt(8)
            run.text = images[1][1] or 'Source: Motherson'
        else:
            shapes._spTree.remove(slide_shapes[11][1]._element)
        
        if len(images) > 2:
            ref_element = slide_shapes[12][1]
            paragraph = ref_element.text_frame.paragraphs[0]
            run = paragraph.runs[0]
            run.text = images[1][1] or 'Source: Motherson'
            run.font.size = Pt(8)
            run.text = images[2][1]
            run.font.color.rgb = RGBColor(255, 255, 255)
        else:
            shapes._spTree.remove(slide_shapes[12][1]._element)

        bars_to_remove = {
            1: [19, 20, 21],
            2: [17, 20, 21],
            3: [17, 19, 21],
            4: [17, 19, 20]
        }

        # Get the list of bar indexes to remove based on tr_level['number']
        bars_to_remove_indexes = bars_to_remove.get(tr_level['number'], [])

        # Remove the bars
        for index in bars_to_remove_indexes:
            shape_to_move = slide_shapes[index][1]
            shapes._spTree.remove(shape_to_move._element)
            
        prod_name = product['name']
        ###***********  insert production code for download here ***********### 

        return factsheet_ppt.save(f'{prod_name}.pptx')  # Save the modified PowerPoint file
    except Exception as e:
            print(str(e))

if __name__ == '__main__':
    generate_pptx()

#https://www.youtube.com/watch?v=SHrR2fFVDO4
